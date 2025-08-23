import os
import sys
import locale
import httpx
from pptx import Presentation
from pptx.util import Pt
from io import BytesIO
from pptx.dml.color import RGBColor
import json
from pathlib import Path
import unicodedata

OPENROUTER_API_KEY = "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"
OPENROUTER_API_URL = "https://api.groq.com/openai/v1/chat/completions"


async def generate_course_content(topic: str, level: str = "débutant", model: str = None) -> dict:
    prompt = f"""
    
"""
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": model if model else "llama-3.1-8b-instant",
        "messages": [
            {"role": "user", "content": prompt}
        ]
    }
    try:
        async with httpx.AsyncClient() as client:
            response = await client.post(OPENROUTER_API_URL, json=data, headers=headers)
            response.raise_for_status()
            return response.json()
    except Exception as e:
        print(f"Error calling OpenRouter API: {e}")
        raise


async def generate_qcm_content(topic: str, level: str = "débutant", model: str = None) -> dict:
    prompt = f"""
Tu es un générateur de QCM. Génère exactement 5 questions à choix multiples sur le thème : "{topic}".
Niveau : {level}.
Structure la réponse exclusivement comme une liste JSON **valide**, sans aucun texte supplémentaire, sans introduction, sans code block.

Format attendu :
[
  {{
    "question": "Quel est ... ?",
    "choices": ["A. ...", "B. ...", "C. ...", "D. ..."],
    "answer": "A"
  }},
  ...
]
"""
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": model if model else "llama-3.1-8b-instant",
        "messages": [
            {"role": "user", "content": prompt}
        ]
    }

    async with httpx.AsyncClient() as client:
        response = await client.post(OPENROUTER_API_URL, headers=headers, json=payload)
        response.raise_for_status()
        result = response.json()
        content = result["choices"][0]["message"]["content"]

        if content.startswith("```"):
            content = content.strip().strip("`").split("json")[-1].strip()

        try:
            parsed_qcm = json.loads(content)
        except json.JSONDecodeError as e:
            raise ValueError(f"Erreur de parsing JSON: {e}\nContenu reçu:\n{content}")

        return {"qcm": parsed_qcm}


async def generate_presentation_slides(topic: str,level: str = "débutant",theorique_hours: int = 0,pratique_hours: int = 0,model: str = None):
    prompt = f"""
    Tu es un assistant expert en création de présentations professionnelles.  
Génère une présentation complète sur le thème : '{topic}'.  
Niveau : {level}.  
Durée indicative : {theorique_hours}h théorique et {pratique_hours}h pratique. 
 
La première diapositive (index 0) doit être un slide de titre :  
- "title" = le titre complet de la présentation (reprenant fidèlement le thème)  
- pas de "content" dans ce slide  

Ensuite :  
- Crée plusieurs diapositives explicatives ("Présentation PPT") avec un titre clair et un ou plusieurs paragraphes développés (pas de puces).  
- Ajoute ensuite des diapositives d’"Exercices pratiques" avec un titre et une consigne détaillée, adaptées au niveau {level}.  

Réponds EXCLUSIVEMENT avec une liste JSON valide, sans texte supplémentaire ni bloc de code.  

Format attendu :  
[  
  {{ "title": "Titre du slide", "content": "Texte explicatif ou consigne" }},  
  ...  
]  """
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": model if model else "llama-3.1-8b-instant",
        "messages": [
            {"role": "user", "content": prompt}
        ]
    }
    try:
        async with httpx.AsyncClient() as client:
            response = await client.post(OPENROUTER_API_URL, json=data, headers=headers)
            response.raise_for_status()
            result = response.json()
            content = result["choices"][0]["message"]["content"]
            if content.startswith("```"):
                content = content.strip().strip("`").split("json")[-1].strip()
            slides = json.loads(content)
            return slides
    except Exception as e:
        print(f"Error calling OpenRouter API: {e}")
        raise


def get_template_path(design_name: str) -> str:
    templates_dir = Path(__file__).parent.parent / "templates"
    templates_dir.mkdir(exist_ok=True)
    
    template_files = {
        "Minimal": "minimal_template.pptx",
        "Corporate": "corporate_template.pptx", 
        "Colorful": "colorful_template.pptx"
    }
    
    template_file = template_files.get(design_name, "minimal_template.pptx")
    template_path = templates_dir / template_file
    
    if not template_path.exists():
        return None
    
    return str(template_path)


def generate_presentation_pptx_with_template(slides, design: str = "Minimal") -> bytes:
    print(f"[DEBUG] sys.getdefaultencoding(): {sys.getdefaultencoding()}")
    print(f"[DEBUG] locale.getpreferredencoding(): {locale.getpreferredencoding()}")
    os.environ["PYTHONIOENCODING"] = "utf-8"

    template_path = get_template_path(design)
    prs = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()

    num_layouts = len(prs.slide_layouts)
    print(f"Layouts: {num_layouts} available.")

    for idx, slide_data in enumerate(slides):
        if idx == 0:
            slide_layout_idx = 0
        elif idx == len(slides) - 1:
            slide_layout_idx = min(17, num_layouts - 1)
        else:
            slide_layout_idx = 4

        slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_idx])

        # --- Title ---
        if slide.shapes.title:
            try:
                title_text = slide_data.get("title", "Titre de la présentation")
                if not isinstance(title_text, str):
                    title_text = str(title_text)
                title_text = unicodedata.normalize("NFC", title_text)
                title_text = title_text.encode("utf-8", errors="replace").decode("utf-8")

                if idx == 0:
                    slide.shapes.title.text = title_text
                elif idx == len(slides) - 1:
                    slide.shapes.title.text = "Résumé et Conclusion"
                else:
                    slide.shapes.title.text = title_text if title_text else f"Slide {idx+1}"
            except Exception as e:
                print(f"[ERROR] Title encoding error: {e}")
                slide.shapes.title.text = "[Erreur d'encodage titre]"

        if idx == 0:
            continue

        content = slide_data.get("content", "").strip()
        if not isinstance(content, str):
            content = str(content)
        content = unicodedata.normalize("NFC", content)
        content_applied = False

        # 1. Try placeholders
        for ph in slide.placeholders:
            if ph.has_text_frame and ph != slide.shapes.title:
                tf = ph.text_frame
                tf.clear()
                for paragraph in content.split("\n"):
                    p = tf.add_paragraph()
                    p.font.size = Pt(24)
                    safe_text = paragraph.strip()
                    if not isinstance(safe_text, str):
                        safe_text = str(safe_text)
                    safe_text = unicodedata.normalize("NFC", safe_text)
                    safe_text = safe_text.encode("utf-8", errors="replace").decode("utf-8")
                    p.text = safe_text
                content_applied = True
                break

        # 2. Try any text box
        if not content_applied:
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    tf = shape.text_frame
                    tf.clear()
                    for paragraph in content.split("\n"):
                        p = tf.add_paragraph()
                        p.font.size = Pt(24)
                        safe_text = paragraph.strip()
                        if not isinstance(safe_text, str):
                            safe_text = str(safe_text)
                        safe_text = unicodedata.normalize("NFC", safe_text)
                        safe_text = safe_text.encode("utf-8", errors="replace").decode("utf-8")
                        p.text = safe_text
                    content_applied = True
                    break

        if not content_applied:
            print(f"⚠ No text frame found for slide {idx+1}")

    pptx_io = BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()


def generate_presentation_pptx(slides) -> bytes:
    prs = Presentation()
    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255) if idx else RGBColor(34, 49, 63)

        slide.shapes.title.text = slide_data.get('title', f"Slide {idx+1}")
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36 if idx else 44)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(34, 49, 63) if idx else RGBColor(255, 255, 255)

        tf = slide.placeholders[1].text_frame
        tf.clear()
        content = slide_data.get('content', '')
        if content:
            for paragraph in content.split('\n'):
                p = tf.add_paragraph()
                safe_text = unicodedata.normalize("NFC", paragraph.strip())
                safe_text = safe_text.encode("utf-8", errors="replace").decode("utf-8")
                p.text = safe_text
                p.font.size = Pt(15)
                p.font.color.rgb = RGBColor(44, 62, 80) if idx else RGBColor(255, 255, 255)

        if tf.paragraphs and not tf.paragraphs[0].text:
            tf._element.remove(tf.paragraphs[0]._p)

    pptx_io = BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()


def generate_course_pptx(title: str, outline: str, summary: str, raw_content: str) -> bytes:
    prs = Presentation()
    # --- Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title or "Titre du cours"
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = summary or "Résumé"
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(34, 49, 63)
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(28)
        slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # --- Outline Slide ---
    outline_slide = prs.slides.add_slide(prs.slide_layouts[1])
    outline_slide.shapes.title.text = "Plan détaillé"
    outline_slide.background.fill.solid()
    outline_slide.background.fill.fore_color.rgb = RGBColor(236, 240, 241)
    outline_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    outline_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    outline_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(34, 49, 63)
    outline_slide.placeholders[1].text = outline or ""
    outline_slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(24)
    outline_slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # --- Summary Slide ---
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_slide.shapes.title.text = "Résumé"
    summary_slide.background.fill.solid()
    summary_slide.background.fill.fore_color.rgb = RGBColor(236, 240, 241)
    summary_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    summary_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    summary_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(34, 49, 63)
    summary_slide.placeholders[1].text = summary or ""
    summary_slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(24)
    summary_slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # --- Main Content Slides ---
    import re
    sections = re.split(r'\n\\s*\\n|(?=^\\d+\\.)', raw_content, flags=re.MULTILINE)
    for section in sections:
        section = section.strip()
        if not section:
            continue
        heading_match = re.match(r'^(\\d+\\.\\s+.+)', section)
        if heading_match:
            heading = heading_match.group(1)
            content = section[len(heading):].strip()
        else:
            heading = "Section"
            content = section

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        slide.shapes.title.text = heading
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(34, 49, 63)

        bullets = [line.strip('-• ') for line in content.split('\n') if line.strip()]
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for bullet in bullets:
            safe_text = unicodedata.normalize("NFC", bullet)
            safe_text = safe_text.encode("utf-8", errors="replace").decode("utf-8")
            p = tf.add_paragraph()
            p.text = safe_text
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(44, 62, 80)

        if tf.paragraphs and not tf.paragraphs[0].text:
            tf._element.remove(tf.paragraphs[0]._p)

    pptx_io = BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()
